#!/usr/bin/env python3

from __future__ import annotations

import hashlib
import re
import shutil
from pathlib import Path
from typing import Iterable

from lxml import etree, html
from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
HTML_PATH = ROOT / "ai_programming_training_zh.html"
PPTX_PATH = ROOT / "ai_programming_training_zh.pptx"
ASSET_DIR = Path("/tmp/ai_training_svg_assets")

FONT_SANS = "PingFang SC"
FONT_MONO = "Menlo"

BG = "F4EFE6"
PAPER = "FFF9F1"
PAPER_SOFT = "FFFBF6"
INK = "1F2620"
MUTED = "5D685E"
LINE = "DED5C7"
ACCENT = "D86F27"
ACCENT_DEEP = "AC4E14"
ACCENT_SOFT = "F7E2D5"
TEAL = "1F6A67"
TEAL_SOFT = "E8F6F5"
OLIVE = "55633B"
OLIVE_SOFT = "EDF3E6"
SLATE_SOFT = "F0F2EE"

SLIDE_W = 13.333
SLIDE_H = 7.5


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color)


def clean_text(text: str | None) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def iter_text(node) -> str:
    if node is None:
        return ""
    return clean_text(" ".join(node.itertext()))


def cls(name: str) -> str:
    return f"contains(concat(' ', normalize-space(@class), ' '), ' {name} ')"


def first(nodes, default=None):
    return nodes[0] if nodes else default


def set_run_style(run, *, font=FONT_SANS, size=12, bold=False, color=INK):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), font)
    r_pr.set(qn("a:cs"), font)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    *,
    font=FONT_SANS,
    size=12,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margins=(0.0, 0.0, 0.0, 0.0),
    autoshrink=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margins[0])
    tf.margin_top = Inches(margins[1])
    tf.margin_right = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    set_run_style(run, font=font, size=size, bold=bold, color=color)
    if autoshrink:
        tf.fit_text(font_family=font, max_size=size, bold=bold)
    return box


def add_text_lines(slide, x, y, w, h, lines, *, margins=(0.0, 0.0, 0.0, 0.0), valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margins[0])
    tf.margin_top = Inches(margins[1])
    tf.margin_right = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = line["text"]
        set_run_style(
            run,
            font=line.get("font", FONT_SANS),
            size=line.get("size", 12),
            bold=line.get("bold", False),
            color=line.get("color", INK),
        )
    return box


def add_round_rect(slide, x, y, w, h, *, fill=PAPER_SOFT, line=LINE, line_w=1.0, radius_shape=SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_w)
    return shape


def add_plain_rect(slide, x, y, w, h, *, fill=PAPER_SOFT, line=LINE, line_w=1.0):
    shape = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_w)
    return shape


def add_oval(slide, x, y, w, h, *, fill, transparency=0.0):
    shape = slide.shapes.add_shape(SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_pill(slide, x, y, text, *, fill=TEAL_SOFT, line=TEAL, font=FONT_SANS, size=9, color=TEAL, h=0.3, padding=0.18):
    w = min(max(0.82, 0.12 * len(text) + padding * 2), 2.3)
    pill = add_round_rect(slide, x, y, w, h, fill=fill, line=line, line_w=0.9)
    add_textbox(
        slide,
        x + 0.02,
        y + 0.01,
        w - 0.04,
        h - 0.02,
        text,
        font=font,
        size=size,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return pill, w


def add_tag_row(slide, tags: Iterable[str], x, y, max_w, *, color_fill=ACCENT_SOFT, color_line=ACCENT, color_text=ACCENT_DEEP, size=8.5):
    cur_x = x
    cur_y = y
    row_h = 0.24
    for tag in tags:
        _, w = add_pill(
            slide,
            cur_x,
            cur_y,
            tag,
            fill=color_fill,
            line=color_line,
            color=color_text,
            size=size,
            h=row_h,
            padding=0.14,
        )
        cur_x += w + 0.06
        if cur_x > x + max_w - 0.8:
            cur_x = x
            cur_y += row_h + 0.05
    return cur_y + row_h


def add_background(slide):
    bg = add_plain_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG, line=BG, line_w=0)
    bg.line.fill.background()
    add_oval(slide, -0.7, -0.4, 2.9, 2.9, fill=TEAL, transparency=0.88)
    add_oval(slide, 10.9, -0.3, 3.0, 3.0, fill=ACCENT, transparency=0.9)
    add_oval(slide, 10.8, 5.6, 2.6, 2.2, fill=OLIVE, transparency=0.92)


def add_panel(slide):
    return add_round_rect(slide, 0.34, 0.28, 12.66, 6.94, fill=PAPER, line=LINE, line_w=1.0)


def add_eyebrow(slide, x, y, text):
    pill = add_round_rect(slide, x, y, 2.95, 0.32, fill=TEAL_SOFT, line=TEAL, line_w=0.9)
    add_textbox(
        slide,
        x + 0.08,
        y + 0.02,
        2.79,
        0.26,
        text,
        size=9,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return pill


def estimate_wrap_width(text: str, base=0.38, factor=0.1, max_w=2.2):
    return min(max(base + factor * len(text), 0.9), max_w)


def add_source_pills(slide, items, x, y, max_w):
    cur_x = x
    cur_y = y
    row_h = 0.3
    for label, href in items:
        w = estimate_wrap_width(label, base=0.42, factor=0.08, max_w=2.1)
        if cur_x + w > x + max_w:
            cur_x = x
            cur_y += row_h + 0.06
        shape = add_round_rect(slide, cur_x, cur_y, w, row_h, fill=TEAL_SOFT, line=TEAL, line_w=0.9)
        box = add_textbox(
            slide,
            cur_x + 0.03,
            cur_y + 0.01,
            w - 0.06,
            row_h - 0.02,
            label,
            size=8.6,
            bold=True,
            color=TEAL,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if href:
            box.text_frame.paragraphs[0].runs[0].hyperlink.address = href
            shape.click_action.hyperlink.address = href
        cur_x += w + 0.08
    return cur_y + row_h


def add_info_chip(slide, x, y, w, h, title, value):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    add_text_lines(
        slide,
        x + 0.14,
        y + 0.12,
        w - 0.28,
        h - 0.24,
        [
            {"text": title, "size": 8.4, "bold": True, "color": MUTED},
            {"text": value, "size": 15.2, "bold": True, "color": INK},
        ],
    )


def add_toc_card(slide, x, y, w, h, num, label):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    add_text_lines(
        slide,
        x + 0.12,
        y + 0.08,
        w - 0.24,
        h - 0.16,
        [
            {"text": num, "size": 8.2, "bold": True, "color": ACCENT_DEEP},
            {"text": label, "size": 10.5, "bold": True, "color": INK},
        ],
    )


def add_note_bar(slide, x, y, w, h, text):
    add_round_rect(slide, x, y, w, h, fill=ACCENT_SOFT, line=ACCENT, line_w=0)
    add_textbox(
        slide,
        x + 0.12,
        y + 0.05,
        w - 0.24,
        h - 0.1,
        text,
        size=10.2,
        color="4A3825",
        bold=False,
        margins=(0.0, 0.0, 0.0, 0.0),
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_section_header(slide, title, kicker, intro, pills):
    add_textbox(slide, 0.66, 0.54, 2.6, 0.2, kicker, size=8.3, bold=True, color=ACCENT_DEEP)
    add_textbox(slide, 0.66, 0.76, 6.0, 0.4, title, size=23, bold=True, color=INK)
    if pills:
        add_source_pills(slide, pills, 8.45, 0.58, 4.0)
    add_line(slide, 0.66, 1.28, 12.6, 1.28, color=LINE, width=0.9)
    add_textbox(slide, 0.66, 1.34, 9.0, 0.28, intro, size=11.5, color=MUTED)


def fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with Image.open(path) as img:
        img_w, img_h = img.size
    scale = min(max_w / img_w, max_h / img_h)
    return img_w * scale, img_h * scale


def add_diagram_panel(slide, x, y, w, h, image_path: Path, caption: str):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    inner_x = x + 0.16
    inner_y = y + 0.14
    inner_w = w - 0.32
    inner_h = h - 0.52
    img_w_px, img_h_px = fit_image(image_path, Inches(inner_w).pt, Inches(inner_h).pt)
    img_w = img_w_px / 72
    img_h = img_h_px / 72
    pic_x = inner_x + (inner_w - img_w) / 2
    pic_y = inner_y + (inner_h - img_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(pic_x), Inches(pic_y), width=Inches(img_w), height=Inches(img_h))
    add_textbox(slide, x + 0.16, y + h - 0.28, w - 0.32, 0.18, caption, size=9.4, color=MUTED)


def add_signal_card(slide, x, y, w, h, card, *, compact=False):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    badge_fill = ACCENT_SOFT
    badge_line = ACCENT
    badge = add_round_rect(slide, x + 0.14, y + 0.14, 0.4, 0.25, fill=badge_fill, line=badge_line, line_w=0.8)
    badge.line.fill.background()
    add_textbox(
        slide,
        x + 0.14,
        y + 0.145,
        0.4,
        0.22,
        card["no"],
        size=8,
        bold=True,
        color=ACCENT_DEEP,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(slide, x + 0.14, y + 0.48, w - 0.28, 0.26, card["title"], size=13.2 if compact else 14.6, bold=True)
    add_textbox(slide, x + 0.14, y + 0.78, w - 0.28, 0.34, card["desc"], size=10.3 if compact else 10.8, color=MUTED)
    if card.get("tags"):
        add_tag_row(slide, card["tags"], x + 0.14, y + h - (0.44 if compact else 0.42), w - 0.28, size=7.8 if compact else 8.2)
    return badge


def add_tool_card(slide, x, y, w, h, card):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    add_textbox(slide, x + 0.12, y + 0.12, w - 0.24, 0.26, card["title"], size=11.7, bold=True)
    if card.get("tags"):
        add_tag_row(slide, card["tags"], x + 0.12, y + 0.44, w - 0.24, size=7.4)
    add_textbox(slide, x + 0.12, y + h - 0.36, w - 0.24, 0.22, card["desc"], size=9.2, color=MUTED)


def add_scenario_card(slide, x, y, w, h, title, desc, *, mono=False, compact=False):
    add_round_rect(slide, x, y, w, h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
    add_textbox(slide, x + 0.12, y + 0.1, w - 0.24, 0.2, title, size=10.2 if compact else 10.7, bold=True)
    add_textbox(
        slide,
        x + 0.12,
        y + 0.33,
        w - 0.24,
        h - 0.4,
        desc,
        size=9.0 if compact else 9.4,
        color=MUTED if not mono else ACCENT_DEEP,
        font=FONT_MONO if mono else FONT_SANS,
        bold=mono,
    )


def add_quote_bar(slide, x, y, w, h, text):
    add_round_rect(slide, x, y, w, h, fill=ACCENT_SOFT, line=ACCENT, line_w=0)
    add_textbox(slide, x + 0.14, y + 0.08, w - 0.28, h - 0.16, text, size=11, bold=True, color="5F3614", valign=MSO_ANCHOR.MIDDLE)


def parse_card(article):
    title = clean_text(" ".join(article.xpath("./h4/text()")))
    desc = clean_text(" ".join(article.xpath("./p//text()")))
    no = clean_text(" ".join(article.xpath(f"./div[{cls('trend-no')}]/text()")))
    tags = [iter_text(span) for span in article.xpath(f".//span[{cls('tag')}]")]
    return {"no": no, "title": title, "desc": desc, "tags": tags}


def extract_diagrams(section):
    diagrams = []
    for wrap in section.xpath(f".//div[{cls('diagram-wrap')}]"):
        svg = first(wrap.xpath(".//*[name()='svg']"))
        if svg is None:
            continue
        diagrams.append(
            {
                "label": svg.get("aria-label"),
                "caption": iter_text(first(wrap.xpath(f".//div[{cls('diagram-caption')}]"))),
            }
        )
    return diagrams


def parse_html():
    doc = html.fromstring(HTML_PATH.read_text(encoding="utf-8"))
    hero = doc.get_element_by_id("top")
    hero_title = first(hero.xpath(".//h1"))
    hero_main = clean_text(" ".join(hero_title.xpath("./text()")))
    hero_highlight = clean_text(" ".join(hero_title.xpath(".//span/text()")))
    data = {
        "hero": {
            "eyebrow": iter_text(first(hero.xpath(f".//div[{cls('eyebrow')}]"))),
            "title_main": hero_main,
            "title_highlight": hero_highlight,
            "lead": iter_text(first(hero.xpath(f".//p[{cls('lead')}]"))),
            "chips": [
                {
                    "k": iter_text(first(node.xpath(f".//span[{cls('k')}]"))),
                    "v": iter_text(first(node.xpath(f".//span[{cls('v')}]"))),
                }
                for node in hero.xpath(f".//div[{cls('chip-row')}]/div[{cls('chip')}]")
            ],
            "toc": [
                {
                    "num": iter_text(first(a.xpath(f".//span[{cls('num')}]"))),
                    "label": iter_text(first(a.xpath(f".//span[{cls('label')}]"))),
                }
                for a in hero.xpath(f".//nav[{cls('toc')}]/a")
            ],
            "note": iter_text(first(hero.xpath(f".//div[{cls('hero-note')}]"))),
            "diagram": extract_diagrams(hero)[0],
        }
    }

    for section_id in ["trends", "tools", "agent", "skills-mcp", "stability", "landing", "sources"]:
        sec = doc.get_element_by_id(section_id)
        data[section_id] = {
            "kicker": iter_text(first(sec.xpath(f".//div[{cls('section-kicker')}]"))),
            "title": iter_text(first(sec.xpath(".//h2"))),
            "intro": iter_text(first(sec.xpath(f".//p[{cls('section-intro')}]"))),
            "sources": [
                (iter_text(a), a.get("href"))
                for a in sec.xpath(f".//div[{cls('source-inline')}]/a")
            ],
            "diagrams": extract_diagrams(sec),
            "signal_cards": [],
        }

    data["trends"]["signal_cards"] = [parse_card(a) for a in data_section(doc, "trends", "signal-grid")]
    data["tools"]["tool_cards"] = [parse_card(a) for a in data_section(doc, "tools", "tool-mini-grid")]
    data["agent"]["scenarios"] = [parse_card(a) for a in data_section(doc, "agent", "scenario-grid")]
    data["agent"]["quote"] = iter_text(first(doc.get_element_by_id("agent").xpath(f".//div[{cls('quote')}]")))
    data["skills-mcp"]["signal_cards"] = [parse_card(a) for a in data_section(doc, "skills-mcp", "signal-grid")]
    data["skills-mcp"]["scenarios"] = [parse_card(a) for a in data_section(doc, "skills-mcp", "scenario-grid")]
    data["stability"]["signal_cards"] = [parse_card(a) for a in data_section(doc, "stability", "signal-grid")]
    data["stability"]["scenarios"] = [parse_card(a) for a in data_section(doc, "stability", "scenario-grid")]
    data["landing"]["signal_cards"] = [parse_card(a) for a in data_section(doc, "landing", "signal-grid")]
    data["sources"]["cloud"] = [
        {
            "title": clean_text(" ".join(a.xpath("./text()"))),
            "desc": iter_text(first(a.xpath("./span"))),
            "href": a.get("href"),
        }
        for a in doc.get_element_by_id("sources").xpath(f".//div[{cls('source-cloud')}]/a")
    ]
    data["sources"]["footnote"] = iter_text(first(doc.get_element_by_id("sources").xpath(f".//p[{cls('footnote')}]")))
    return doc, data


def data_section(doc, section_id, grid_class):
    grid = first(doc.get_element_by_id(section_id).xpath(f".//div[{cls(grid_class)}]"))
    return list(grid.xpath("./article")) if grid is not None else []


def export_svgs(doc) -> dict[str, Path]:
    if ASSET_DIR.exists():
        shutil.rmtree(ASSET_DIR)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    images = {}
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        for svg in doc.xpath("//*[name()='svg']"):
            label = svg.get("aria-label")
            if not label:
                continue
            svg_copy = etree.fromstring(etree.tostring(svg))
            if "xmlns" not in svg_copy.attrib:
                svg_copy.set("xmlns", "http://www.w3.org/2000/svg")
            view_box = svg_copy.get("viewBox")
            width = 960
            height = 640
            if view_box:
                parts = [float(part) for part in re.split(r"[,\s]+", view_box.strip()) if part]
                if len(parts) == 4 and parts[2] > 0 and parts[3] > 0:
                    width = parts[2]
                    height = parts[3]
            scale = 3
            render_w = max(1200, int(width * scale))
            render_h = max(700, int(height * scale))
            svg_copy.set("width", str(render_w))
            svg_copy.set("height", str(render_h))
            digest = hashlib.md5(label.encode("utf-8")).hexdigest()[:10]
            svg_text = etree.tostring(svg_copy, encoding="unicode")
            png_path = ASSET_DIR / f"{digest}.png"
            page = browser.new_page(viewport={"width": render_w, "height": render_h}, device_scale_factor=1)
            page.set_content(
                (
                    "<html><head><style>"
                    "html,body{margin:0;padding:0;background:transparent;}"
                    "svg,text{font-family:'PingFang SC','Noto Sans SC',sans-serif;}"
                    "</style></head><body>"
                    f"{svg_text}"
                    "</body></html>"
                )
            )
            page.screenshot(path=str(png_path), omit_background=True)
            page.close()
            images[label] = png_path
        browser.close()
    return images


def build_hero(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_eyebrow(slide, 0.7, 0.58, data["eyebrow"])
    add_textbox(slide, 0.8, 1.02, 5.6, 0.55, data["title_main"], size=28, bold=True)
    add_textbox(slide, 0.8, 1.56, 5.9, 0.58, data["title_highlight"], size=31, bold=True, color=ACCENT_DEEP)
    add_textbox(slide, 0.8, 2.18, 5.3, 0.24, data["lead"], size=12.5, color=MUTED)

    chip_w = 2.8
    chip_h = 0.88
    chip_x = [0.8, 3.76]
    chip_y = [2.66, 3.66]
    for idx, chip in enumerate(data["chips"]):
        add_info_chip(slide, chip_x[idx % 2], chip_y[idx // 2], chip_w, chip_h, chip["k"], chip["v"])

    toc_w = 1.78
    toc_h = 0.62
    toc_start_x = 0.8
    toc_start_y = 4.82
    toc_gap_x = 0.16
    toc_gap_y = 0.14
    for idx, item in enumerate(data["toc"]):
        row, col = divmod(idx, 3)
        x = toc_start_x + col * (toc_w + toc_gap_x)
        y = toc_start_y + row * (toc_h + toc_gap_y)
        add_toc_card(slide, x, y, toc_w, toc_h, item["num"], item["label"])

    add_note_bar(slide, 0.8, 6.28, 5.84, 0.56, data["note"])

    diagram = data["diagram"]
    add_diagram_panel(slide, 7.03, 0.72, 5.42, 5.92, images[diagram["label"]], diagram["caption"])


def build_trends(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.76, 7.08, 4.95, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    card_w = 2.13
    card_h = 2.32
    start_x = 7.94
    start_y = 1.76
    for idx, card in enumerate(data["signal_cards"]):
        row, col = divmod(idx, 2)
        add_signal_card(slide, start_x + col * (card_w + 0.14), start_y + row * (card_h + 0.16), card_w, card_h, card)


def build_tools(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.76, 6.82, 4.95, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    card_w = 2.18
    card_h = 1.49
    start_x = 7.66
    start_y = 1.76
    for idx, card in enumerate(data["tool_cards"]):
        row, col = divmod(idx, 2)
        add_tool_card(slide, start_x + col * (card_w + 0.16), start_y + row * (card_h + 0.15), card_w, card_h, card)


def build_agent(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.76, 5.93, 3.88, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    add_diagram_panel(slide, 6.74, 1.76, 5.93, 3.88, images[data["diagrams"][1]["label"]], data["diagrams"][1]["caption"])
    scenario_w = 2.9
    scenario_h = 0.72
    for idx, card in enumerate(data["scenarios"]):
        add_scenario_card(slide, 0.66 + idx * (scenario_w + 0.13), 5.8, scenario_w, scenario_h, card["title"], card["desc"], compact=True)
    add_quote_bar(slide, 0.66, 6.64, 12.01, 0.4, data["quote"])


def build_skill_mcp(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.74, 5.93, 3.5, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    add_diagram_panel(slide, 6.74, 1.74, 5.93, 3.5, images[data["diagrams"][1]["label"]], data["diagrams"][1]["caption"])
    card_w = 2.9
    signal_y = 5.38
    for idx, card in enumerate(data["signal_cards"]):
        add_signal_card(slide, 0.66 + idx * (card_w + 0.13), signal_y, card_w, 0.72, card, compact=True)
    for idx, card in enumerate(data["scenarios"]):
        add_scenario_card(
            slide,
            0.66 + idx * (card_w + 0.13),
            6.16,
            card_w,
            0.62,
            card["title"],
            card["desc"],
            mono="@" in card["desc"] or "use" in card["desc"],
            compact=True,
        )


def build_stability(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.74, 5.93, 3.5, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    add_diagram_panel(slide, 6.74, 1.74, 5.93, 3.5, images[data["diagrams"][1]["label"]], data["diagrams"][1]["caption"])
    card_w = 2.9
    for idx, card in enumerate(data["signal_cards"]):
        add_signal_card(slide, 0.66 + idx * (card_w + 0.13), 5.38, card_w, 0.72, card, compact=True)
    for idx, card in enumerate(data["scenarios"]):
        add_scenario_card(slide, 0.66 + idx * (card_w + 0.13), 6.16, card_w, 0.62, card["title"], card["desc"], compact=True)


def build_landing(prs, data, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], data["intro"], data["sources"])
    add_diagram_panel(slide, 0.66, 1.76, 7.08, 4.95, images[data["diagrams"][0]["label"]], data["diagrams"][0]["caption"])
    card_w = 2.13
    card_h = 2.22
    start_x = 7.94
    start_y = 1.96
    for idx, card in enumerate(data["signal_cards"]):
        row, col = divmod(idx, 2)
        add_signal_card(slide, start_x + col * (card_w + 0.14), start_y + row * (card_h + 0.16), card_w, card_h, card)


def build_sources(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_panel(slide)
    add_section_header(slide, data["title"], data["kicker"], "", [])
    card_w = 3.92
    card_h = 0.68
    start_x = 0.66
    start_y = 1.58
    gap_x = 0.16
    gap_y = 0.12
    for idx, item in enumerate(data["cloud"]):
        row, col = divmod(idx, 3)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_round_rect(slide, x, y, card_w, card_h, fill=PAPER_SOFT, line=LINE, line_w=0.9)
        box = add_text_lines(
            slide,
            x + 0.12,
            y + 0.1,
            card_w - 0.24,
            card_h - 0.2,
            [
                {"text": item["title"], "size": 9.2, "bold": True, "color": TEAL},
                {"text": item["desc"], "size": 7.8, "color": MUTED},
            ],
        )
        if item["href"]:
            box.text_frame.paragraphs[0].runs[0].hyperlink.address = item["href"]
    add_textbox(slide, 0.66, 6.9, 8.0, 0.2, data["footnote"], size=8.8, color=MUTED)


def main():
    doc, data = parse_html()
    images = export_svgs(doc)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_hero(prs, data["hero"], images)
    build_trends(prs, data["trends"], images)
    build_tools(prs, data["tools"], images)
    build_agent(prs, data["agent"], images)
    build_skill_mcp(prs, data["skills-mcp"], images)
    build_stability(prs, data["stability"], images)
    build_landing(prs, data["landing"], images)
    build_sources(prs, data["sources"])
    prs.save(PPTX_PATH)
    print(f"Wrote {PPTX_PATH}")


if __name__ == "__main__":
    main()
